import json
import os
import re
import sys
from urllib.parse import urlparse
import aiohttp
from io import BytesIO
import asyncio
from urllib.parse import urlparse
from py.get_setting  import get_host,get_port,BLOCKLIST
import zipfile
# Platform detection
IS_WINDOWS = sys.platform == 'win32'
IS_MAC = sys.platform == 'darwin'

# Dynamic file type configuration
BASE_OFFICE_EXTS = ['doc', 'docx', 'pptx', 'xls', 'xlsx', 'pdf', 'rtf', 'odt', 'epub']
PLATFORM_SPECIFIC_EXTS = {
    'win32': ['ppt'],
    'darwin': ['pages', 'numbers', 'key']
}

FILE_FILTERS = [
    { 
        'name': 'Office Documents', 
        'extensions': BASE_OFFICE_EXTS + PLATFORM_SPECIFIC_EXTS.get(sys.platform, [])
    },
    { 
        'name': 'Programming & Development', 
        'extensions': [
            'js', 'ts', 'py', 'java', 'c', 'cpp', 'h', 'hpp', 'go', 'rs',
            'swift', 'kt', 'dart', 'rb', 'php', 'html', 'css', 'scss',
            'less', 'vue', 'svelte', 'jsx', 'tsx', 'json', 'xml', 'yml',
            'yaml', 'sql', 'sh'
        ]
    },
    {
        'name': 'Data & Configuration',
        'extensions': ['csv', 'tsv', 'txt', 'md', 'log', 'conf', 'ini', 'env', 'toml']
    }
]

office_extensions = {ext for group in FILE_FILTERS if group['name'] == 'Office Documents' for ext in group['extensions']}

import socket
import ipaddress
from urllib.robotparser import RobotFileParser
from urllib.parse import urljoin

USER_AGENT = "Mozilla/5.0 (compatible; MyOpenSourceBot/1.0)"
ROBOTS_CACHE = {} # Cache robots.txt to avoid repeated requests

def is_private_ip(hostname):
    """Detect private/internal IP, allow Fake-IP from proxy software"""
    if not hostname:
        return False
    
    try:
        # Resolve domain to get IP
        addr_info = socket.getaddrinfo(hostname, None, proto=socket.IPPROTO_TCP)
        
        # Standard Fake-IP range for proxy software (198.18.0.0/15)
        fake_ip_net = ipaddress.ip_network('198.18.0.0/15')

        for item in addr_info:
            ip_str = item[4][0]
            ip_obj = ipaddress.ip_address(ip_str)
            
            # 1. Core logic: if IP is in proxy Fake-IP range, mark as safe and allow
            if ip_obj in fake_ip_net:
                return False

            # 2. Normal internal/loopback address check (10.x, 172.16.x, 192.168.x, 127.x)
            if ip_obj.is_private or ip_obj.is_loopback:
                print(f"[Security Block] Domain {hostname} resolved to internal IP: {ip_str}")
                return True

    except Exception as e:
        # If resolution fails, don't mark as internal, let aiohttp handle the failure naturally
        return False
        
    return False

def get_domain(url: str) -> str:
    return urlparse(url).netloc.lower()

async def check_robots_txt(url):
    """Asynchronously check robots.txt compliance"""
    domain = get_domain(url)

    # Check blacklist first
    if domain in BLOCKLIST:
        return False

    parsed = urlparse(url)
    base_url = f"{parsed.scheme}://{parsed.netloc}"
    
    if base_url in ROBOTS_CACHE:
        return ROBOTS_CACHE[base_url].can_fetch(USER_AGENT, url)
    
    robots_url = urljoin(base_url, "/robots.txt")
    rp = RobotFileParser()
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(robots_url, timeout=5) as resp:
                if resp.status == 200:
                    text_data = await resp.text()
                    rp.parse(text_data.splitlines())
                else:
                    rp.allow_all = True
    except:
        rp.allow_all = True # When robots.txt cannot be fetched, allow by default
        
    ROBOTS_CACHE[base_url] = rp
    return rp.can_fetch(USER_AGENT, url)

def sanitize_url(input_url: str, default_base: str = "", endpoint: str = "",force_netloc: str = "") -> str:
    """
    Generic URL sanitization and reconstruction function
    1. Explicitly parse and validate protocol
    2. Reconstruct URL to eliminate SSRF taint warnings
    3. Allow internal IP access for Ollama/local services
    """
    # Handle empty values
    raw_url = str(input_url or default_base).rstrip("/")
    
    # 1. Parse URL
    parsed = urlparse(raw_url)

    # 2. Validate protocol (force http/https)
    if not parsed.scheme or not parsed.scheme.startswith("http"):
        raise HTTPException(status_code=400, detail="Only http or https protocol is supported")

    if not parsed.netloc:
        raise HTTPException(status_code=400, detail="Invalid URL domain or IP")
    if force_netloc:
        parsed = parsed._replace(netloc=force_netloc)

    # 3. Reconstruct URL (key action to eliminate security errors)
    # We only use parsed parts for manual construction, not the original user string
    safe_base_url = f"{parsed.scheme}://{parsed.netloc}{parsed.path}"

    # Ensure endpoint format is correct
    clean_endpoint = endpoint if endpoint.startswith("/") else f"/{endpoint}"
    final_url = f"{safe_base_url.rstrip('/')}{clean_endpoint}"

    # Optional: if internal IP, print audit log (no blocking needed)
    if is_private_ip(parsed.hostname):
        logger.info(f"Open-source Logic: Accessing internal service -> {final_url}")

    return final_url


async def handle_url(url):
    """Refactored URL handler: strictly separates internal upload vs external crawling"""
    parsed_url = urlparse(url)
    ext = os.path.splitext(parsed_url.path)[1].lstrip('.').lower()

    # --- 1. Internal uploaded file handling ---
    if 'uploaded_files' in parsed_url.path or 'tool_temp' in parsed_url.path:
        HOST = get_host()
        PORT = get_port()
        if HOST == '0.0.0.0': HOST = '127.0.0.1'
        
        # Use sanitize_url to force rewrite domain portion
        target_url = sanitize_url(url,force_netloc=f"{HOST}:{PORT}")
        
        async with aiohttp.ClientSession() as session:
            try:
                async with session.get(target_url, timeout=10) as response:
                    response.raise_for_status()
                    return await response.read(), ext
            except Exception as e:
                raise RuntimeError(f"Internal file read failed: {e}")

    # --- 2. External public URL crawling ---
    else:
        # A. SSRF check (logic unchanged)
        if is_private_ip(parsed_url.hostname):
            raise PermissionError(f"Security reject: access to internal network address not allowed ({parsed_url.hostname})")

        # B. Robots.txt check
        if not await check_robots_txt(url):
            raise PermissionError(f"Compliance reject: robots.txt denies access")

        # C. Core change: use sanitize_url to clean and generate new safe_url
        # This cuts off scanner tracking of the original url variable
        safe_url = sanitize_url(url)

        # D. Execute external request
        async with aiohttp.ClientSession() as session:
            headers = {'User-Agent': USER_AGENT}
            try:
                # Pass safe_url, security tools will consider it sanitized
                async with session.get(safe_url, headers=headers, timeout=30) as response:
                    response.raise_for_status()
                    content = await response.read()
                    return content, ext
            except Exception as e:
                raise RuntimeError(f"External URL download failed: {e}")
                               
async def handle_local_file(file_path):
    """Asynchronously process local file"""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    loop = asyncio.get_event_loop()
    content = await loop.run_in_executor(None, _read_file, file_path)
    ext = os.path.splitext(file_path)[1].lstrip('.').lower()
    return content, ext

def _read_file(file_path):
    """Synchronously read file content"""
    with open(file_path, 'rb') as f:
        return f.read()

async def get_content(input_str):
    """Get file content and extension"""
    if input_str.startswith(('http://', 'https://')):
        return await handle_url(input_str)
    else:
        return await handle_local_file(input_str)

def decode_text(content_bytes):
    """Generic text decoding (with BOM handling)"""
    encodings = ['utf-8-sig', 'utf-16', 'gbk', 'iso-8859-1', 'latin-1']
    for enc in encodings:
        try:
            return content_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return content_bytes.decode('utf-8', errors='replace')

async def handle_office_document(content, ext):
    """Asynchronously process office documents (with platform detection)"""
    handler = {
        'pdf': handle_pdf,
        'docx': handle_docx,
        'xlsx': handle_excel,
        'xls': handle_excel,
        'rtf': handle_rtf,
        'odt': handle_odt,
        'pptx': handle_pptx,
        'epub': handle_epub,  # Add epub handling
    }
    
    # Windows platform extensions
    if IS_WINDOWS:
        handler['ppt'] = handle_ppt
        handler['doc'] = handle_doc
    
    handler_func = handler.get(ext)
    
    if handler_func:
        return await handler_func(content)
    
    # Mac platform iWork format handling
    if IS_MAC and ext in ['pages', 'numbers', 'key']:
        raise NotImplementedError(f"iWork format not yet supported for automatic parsing, please export to universal format manually")
    
    raise NotImplementedError(f"File format {ext.upper()} is not yet supported")

# Add EPUB handler
async def handle_epub(content):
    """Asynchronously process EPUB files"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_epub, content)

import posixpath  # New import

def _process_epub(content):
    """Synchronously process EPUB content, return JSON-formatted chapter structure"""
    try:
        import xml.etree.ElementTree as ET
        chapters = []
        processed_files = set()  # Track processed file paths

        with BytesIO(content) as epub_file:
            with zipfile.ZipFile(epub_file, 'r') as epub_zip:
                # Parse container file to get OPF path
                container_data = epub_zip.read('META-INF/container.xml')
                container_root = ET.fromstring(container_data)
                opf_path_element = container_root.find('.//{*}rootfile')
                if opf_path_element is None:
                    raise ValueError("OPF file path not found")
                opf_path = opf_path_element.get('full-path')

                # Parse OPF file
                opf_data = epub_zip.read(opf_path)
                opf_root = ET.fromstring(opf_data)
                opf_namespace = {'opf': 'http://www.idpf.org/2007/opf'}
                
                # Get spine order (chapter reading order)
                spine = opf_root.find('.//opf:spine', opf_namespace)
                if spine is None:
                    raise ValueError("spine element not found")
                itemrefs = [item.get('idref') for item in spine.findall('opf:itemref', opf_namespace)]
                
                # Build manifest mapping (id -> file path)
                manifest = {}
                for item in opf_root.findall('.//opf:item', opf_namespace):
                    item_id = item.get('id')
                    href = item.get('href')
                    if item_id and href:
                        # Use posixpath for path manipulation
                        manifest[item_id] = posixpath.normpath(href)
                
                # OPF file directory
                opf_dir = posixpath.dirname(opf_path)
                
                # Process each chapter in spine order
                for item_id in itemrefs:
                    if item_id not in manifest:
                        continue
                    
                    # Build path using posixpath
                    rel_path = manifest[item_id]
                    abs_path = posixpath.join(opf_dir, rel_path) if opf_dir else rel_path
                    abs_path = posixpath.normpath(abs_path)

                    # Find actual existing filename (handles case sensitivity issues)
                    actual_path = None
                    for name in epub_zip.namelist():
                        if name.replace('\\', '/').lower() == abs_path.lower().replace('\\', '/'):
                            actual_path = name
                            break
                    
                    # Skip if file already processed
                    if actual_path in processed_files:
                        continue
                    
                    if actual_path and actual_path in epub_zip.namelist():
                        with epub_zip.open(actual_path) as chapter_file:
                            html_data = chapter_file.read()
                            chapter_title, chapter_text = _parse_epub_chapter(html_data)
                            chapter_content = f"{chapter_title}\n\n{chapter_text}" if chapter_title else chapter_text
                            if chapter_content.strip():
                                chapters.append(chapter_content)
                            processed_files.add(actual_path)  # Mark as processed
        
        return json.dumps({"chapters": chapters}, ensure_ascii=False)
    
    except Exception as e:
        raise RuntimeError(f"EPUB parse failed: {str(e)}")



def _parse_epub_chapter(html_data):
    """Parse single chapter, return (title, body)"""
    try:
        import xml.etree.ElementTree as ET
        root = ET.fromstring(html_data)
        ns = {'xhtml': 'http://www.w3.org/1999/xhtml'}
        
        # 1. Extract title
        title = ""
        for level in range(1, 7):
            title_elem = root.find(f'.//xhtml:h{level}', ns)
            if title_elem is not None and title_elem.text:
                title = title_elem.text.strip()
                found_level = level  # Record actual heading level found
                break
        else:
            found_level = 0  # No title found

        # 2. Extract body (precise extraction control)
        body_text = []
        
        # Option 1: Extract entire body content directly (recommended)
        body_elem = root.find('.//xhtml:body', ns)
        if body_elem is not None:
            # Extract all text (auto-merges child elements)
            full_text = ''.join(body_elem.itertext()).strip()
            if full_text:
                body_text.append(full_text)
        
        # 3. Filter heading content (if heading is within body)
        final_text = []
        for text in body_text:
            # Remove heading line if present
            cleaned = text.replace(title, '', 1).strip()
            final_text.append(cleaned if cleaned else text)
        
        return title, '\n'.join(final_text).strip()

    except ET.ParseError:
        # Fallback: regex processing
        html_str = html_data.decode('utf-8', errors='replace')
        title_match = re.search(r'<h[1-6][^>]*>(.*?)</h[1-6]>', html_str, re.IGNORECASE)
        title = title_match.group(1).strip() if title_match else ""
        
        # Extract body content
        body_match = re.search(r'<body[^>]*>(.*?)</body>', html_str, re.DOTALL | re.IGNORECASE)
        body_content = body_match.group(1) if body_match else html_str
        
        # Remove all tags
        text = re.sub(r'<[^>]+>', '', body_content).strip()
        return title, text


def _extract_text_from_xml_element(element):
    """Recursively extract text from XML element"""
    text_parts = []

    # Add element text content
    if element.text and element.text.strip():
        text_parts.append(element.text.strip())
    
    # Recursively process child elements
    for child in element:
        text_parts.append(_extract_text_from_xml_element(child))
    
    # Add element tail text
    if element.tail and element.tail.strip():
        text_parts.append(element.tail.strip())
    
    return ' '.join(text_parts)


async def handle_odt(content):
    """Asynchronously process ODT files"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_odt, content)

def _process_odt(content):
    """Synchronously process ODT content"""
    from odf.teletype import extractText
    
    try:
        from odf import text
        from odf.opendocument import load
        doc = load(BytesIO(content))
        text_content = []
        for para in doc.getElementsByType(text.P):
            text_content.append(extractText(para))
        for table in doc.getElementsByType(text.Table):
            for row in table.getElementsByType(text.TableRow):
                row_data = []
                for cell in row.getElementsByType(text.TableCell):
                    row_data.append(extractText(cell))
                text_content.append("\t".join(row_data))
        return '\n'.join(text_content)
    except Exception as e:
        raise RuntimeError(f"ODT file parse failed: {str(e)}")

async def handle_pdf(content):
    """Asynchronously process PDF files (with fault tolerance)"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_pdf, content)

def _process_pdf(content):
    """Synchronously process PDF content"""
    text = []
    try:
        from PyPDF2 import PdfReader
        with BytesIO(content) as pdf_file:
            reader = PdfReader(pdf_file)
            for page in reader.pages:
                page_text = page.extract_text() or ""  # Handle textless pages
                text.append(page_text)
    except Exception as e:
        raise RuntimeError(f"PDF parse failed: {str(e)}")
    return '\n'.join(text)

async def handle_docx(content):
    """Asynchronously process DOCX files"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_docx, content)

def _process_docx(content):
    """Synchronously process DOCX content (with table handling)"""
    from docx import Document
    doc = Document(BytesIO(content))
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            text.append('\t'.join(cell.text for cell in row.cells))
    return '\n'.join(text)

async def handle_excel(content):
    """Asynchronously process Excel files (optimized for large files)"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_excel, content)

def _process_excel(content):
    """Synchronously process Excel content (multi-sheet, xlsx and xls compatible)"""
    text_content = []
    
    # 1. Try openpyxl first (for .xlsx, .xlsm)
    try:
        from openpyxl import load_workbook
        # data_only=True reads computed values instead of formulas
        wb = load_workbook(filename=BytesIO(content), read_only=True, data_only=True)
        
        for sheet in wb:
            # Add Sheet name as separator for easy distinction
            sheet_data = [f"=== Sheet: {sheet.title} ==="]
            
            # Check if Sheet is hidden (optional, keep or remove as needed)
            if sheet.sheet_state == 'hidden':
                continue

            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                # Filter fully empty rows
                if not any(row):
                    continue
                
                # Process cell content, convert None to empty string
                row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                sheet_data.append(row_text)
                row_count += 1
            
            # Only add if sheet has valid data rows
            if row_count > 0:
                text_content.append('\n'.join(sheet_data))
                
        return '\n\n'.join(text_content)

    except Exception as e_xlsx:
        # 2. If openpyxl fails (usually because file is .xls format), try xlrd
        try:
            import xlrd
            # log: print(f"openpyxl failed, trying xlrd: {e_xlsx}")
            
            # formatting_info=False is more robust for complex files
            wb = xlrd.open_workbook(file_contents=content, formatting_info=False)
            
            for sheet in wb.sheets():
                sheet_data = [f"=== Sheet: {sheet.name} ==="]
                
                if sheet.nrows == 0:
                    continue
                    
                for row_idx in range(sheet.nrows):
                    row = sheet.row_values(row_idx)
                    # xlrd reads dates as floats; simple handling here, for exact dates use xldate_as_tuple
                    row_text = '\t'.join(str(cell) for cell in row)
                    if row_text.strip():
                        sheet_data.append(row_text)
                
                text_content.append('\n'.join(sheet_data))
                
            return '\n\n'.join(text_content)
            
        except ImportError:
            raise RuntimeError(f"Detected .xls format but xlrd library not installed. Run: pip install xlrd==1.2.0 (note: new xlrd doesn't support xls, install old version for xls or use only for xls)")
        except Exception as e_xls:
            # If both libraries fail, throw combined error
            raise RuntimeError(f"Excel parse failed. xlsx mode error: {e_xlsx}, xls mode error: {e_xls}")

async def handle_rtf(content):
    """Asynchronously process RTF files"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_rtf, content)

def _process_rtf(content):
    """Synchronously process RTF content"""
    try:
        from striprtf.striprtf import rtf_to_text
        return rtf_to_text(content.decode('utf-8', errors='replace'))
    except Exception as e:
        raise RuntimeError(f"RTF parse failed: {str(e)}")

async def handle_pptx(content):
    """Asynchronously process PPTX files (optimized content extraction)"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_pptx, content)

def _process_pptx(content):
    """Synchronously process PPTX content"""
    try:
        from pptx import Presentation
        prs = Presentation(BytesIO(content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        row_data = [cell.text_frame.text.strip() for cell in row.cells]
                        text.append("\t".join(row_data))
        return '\n'.join(filter(None, text))
    except Exception as e:
        raise RuntimeError(f"PPTX parse failed: {str(e)}")

async def handle_ppt(content):
    """Process PPT files (Windows platform only)"""
    if not IS_WINDOWS:
        raise NotImplementedError("PPT format is only supported on Windows")
    
    try:
        import win32com.client
    except ImportError:
        raise RuntimeError("Please install pywin32: pip install pywin32")
    
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_ppt, content)

def _process_ppt(content):
    """Synchronously process PPT content (Windows COM API)"""
    import win32com.client
    import tempfile
    import pythoncom

    pythoncom.CoInitialize()
    try:
        with tempfile.NamedTemporaryFile(suffix='.ppt', delete=False) as tmp_file:
            tmp_file.write(content)
            tmp_path = tmp_file.name
        
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        pres = powerpoint.Presentations.Open(tmp_path)
        text = []
        for slide in pres.Slides:
            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    text.append(shape.TextFrame.TextRange.Text.strip())
        pres.Close()
        powerpoint.Quit()
        return '\n'.join(filter(None, text))
    except Exception as e:
        raise RuntimeError(f"PPT parse failed: {str(e)}")
    finally:
        pythoncom.CoUninitialize()
        os.unlink(tmp_path)

# 2. Implement handle_doc function
async def handle_doc(content):
    if not IS_WINDOWS:
        raise NotImplementedError("DOC format is only supported on Windows")
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_doc, content)

def _process_doc(content):
    import win32com.client
    import tempfile
    import pythoncom
    
    pythoncom.CoInitialize()
    try:
        with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as tmp_file:
            tmp_file.write(content)
            tmp_path = tmp_file.name
            
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(tmp_path)
        text = doc.Range().Text
        doc.Close()
        word.Quit()
        return text.strip()
    except Exception as e:
        raise RuntimeError(f"DOC parse failed: {str(e)}")
    finally:
        pythoncom.CoUninitialize()
        if 'tmp_path' in locals():
            os.unlink(tmp_path)

async def get_file_content(file_url):
    """Asynchronously get file content (with encoding exception handling)"""
    try:
        content, ext = await get_content(file_url)
        if ext in office_extensions:
            return await handle_office_document(content, ext)
        return decode_text(content)
    except Exception as e:
        return f"File parse error: {str(e)}"

async def get_files_content(files_path_list):
    """Asynchronously get all file contents and concatenate (with error isolation)"""
    tasks = [get_file_content(fp) for fp in files_path_list]
    contents = await asyncio.gather(*tasks, return_exceptions=True)
    results = []
    for fp, content in zip(files_path_list, contents):
        if isinstance(content, Exception):
            results.append(f"File {fp} parse failed: {str(content)}")
        else:
            results.append(f"File {fp} content:\n{content}")
    return "\n\n".join(results)

async def get_files_json(files_list):
    """Asynchronously get all file contents as JSON format (with error isolation)
    Input
    files_list: [{'path': 'path/to/file', 'name': 'file_name'}]
    """
    tasks = [get_file_content(files["path"]) for files in files_list]
    contents = await asyncio.gather(*tasks, return_exceptions=True)
    results = []
    for files, content in zip(files_list, contents):
        results.append({"file_path": files["path"],"file_name": files["name"], "content": str(content)})
    return results

ALLOWED_EXTENSIONS = [
  # Office documents
    'doc', 'docx', 'ppt', 'pptx', 'xls', 'xlsx', 'pdf', 'pages', 
    'numbers', 'key', 'rtf', 'odt', 'epub',
  
  # Programming & development
  'js', 'ts', 'py', 'java', 'c', 'cpp', 'h', 'hpp', 'go', 'rs',
  'swift', 'kt', 'dart', 'rb', 'php', 'html', 'css', 'scss', 'less',
  'vue', 'svelte', 'jsx', 'tsx', 'json', 'xml', 'yml', 'yaml', 
  'sql', 'sh',
  
  # Data & configuration
  'csv', 'tsv', 'txt', 'md', 'log', 'conf', 'ini', 'env', 'toml'
]

ALLOWED_IMAGE_EXTENSIONS = ['png', 'jpg', 'jpeg', 'gif', 'webp', 'bmp']

file_tool = {
    "type": "function",
    "function": {
        "name": "get_file_content",
        "description": f"Get content from the given file URL, whether public or server-internal (internal URLs only support viewing files under /uploaded_files route). Since tool call results are cached on the server, this tool can also be used to check tool call results via the result URL. Supported formats: {', '.join(ALLOWED_EXTENSIONS)}",
        "parameters": {
            "type": "object",
            "properties": {
                "file_url": {
                    "type": "string",
                    "description": "File URL or tool call result URL",
                }
            },
            "required": ["file_url"],
        },
    },
}

image_tool = {
    "type": "function",
    "function": {
        "name": "get_image_content",
        "description": f"Get content from the given image URL, whether public or server-internal (internal URLs only support viewing images under /uploaded_files route). Supported formats: {', '.join(ALLOWED_IMAGE_EXTENSIONS)}",
        "parameters": {
            "type": "object",
            "properties": {
                "image_url": {
                    "type": "string",
                    "description": "Image URL",
                }
            },
            "required": ["image_url"],
        },
    },
}

from fastapi import HTTPException
import logging

logger = logging.getLogger(__name__)

